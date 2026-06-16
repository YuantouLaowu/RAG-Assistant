import os
from typing import List, Dict, Optional

import docx2txt
from PyPDF2 import PdfReader
from pptx import Presentation

from config import DATA_DIR


class DocumentLoader:
    def __init__(
        self,
        data_dir: str = DATA_DIR,
    ):
        self.data_dir = data_dir
        self.supported_formats = [".pdf", ".pptx", ".docx", ".txt"]

    def load_pdf(self, file_path: str) -> List[Dict]:
        """加载PDF文件，按页返回内容"""
        pdf_content = []
        try:
            reader = PdfReader(file_path)
            for i, page in enumerate(reader.pages):
                text = page.extract_text() or ""
                formatted_text = f"--- 第 {i+1} 页 ---\n{text}\n"
                pdf_content.append({"text": formatted_text})
        except Exception as e:
            print(f"Error loading PDF {file_path}: {e}")
        return pdf_content

    def load_pptx(self, file_path: str) -> List[Dict]:
        """加载PPT文件，按幻灯片返回内容"""
        pptx_content = []
        try:
            prs = Presentation(file_path)
            for i, slide in enumerate(prs.slides):
                text_runs = []
                for shape in slide.shapes:
                    if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                        text_runs.append(shape.text)
                
                text = "\n".join(text_runs)
                formatted_text = f"--- 幻灯片 {i+1} ---\n{text}\n"
                pptx_content.append({"text": formatted_text})
        except Exception as e:
            print(f"Error loading PPTX {file_path}: {e}")
        return pptx_content

    def load_docx(self, file_path: str) -> str:
        """加载DOCX文件"""
        try:
            return docx2txt.process(file_path)
        except Exception as e:
            print(f"Error loading DOCX {file_path}: {e}")
            return ""

    def load_txt(self, file_path: str) -> str:
        """加载TXT文件"""
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                return f.read()
        except Exception as e:
            print(f"Error loading TXT {file_path}: {e}")
            return ""

    def load_document(self, file_path: str) -> List[Dict[str, str]]:
        """加载单个文档，PDF和PPT按页/幻灯片分割，返回文档块列表"""
        ext = os.path.splitext(file_path)[1].lower()
        filename = os.path.basename(file_path)
        documents = []

        if ext == ".pdf":
            pages = self.load_pdf(file_path)
            for page_idx, page_data in enumerate(pages, 1):
                documents.append(
                    {
                        "content": page_data["text"],
                        "filename": filename,
                        "filepath": file_path,
                        "filetype": ext,
                        "page_number": page_idx,
                    }
                )
        elif ext == ".pptx":
            slides = self.load_pptx(file_path)
            for slide_idx, slide_data in enumerate(slides, 1):
                documents.append(
                    {
                        "content": slide_data["text"],
                        "filename": filename,
                        "filepath": file_path,
                        "filetype": ext,
                        "page_number": slide_idx,
                    }
                )
        elif ext == ".docx":
            content = self.load_docx(file_path)
            if content:
                documents.append(
                    {
                        "content": content,
                        "filename": filename,
                        "filepath": file_path,
                        "filetype": ext,
                        "page_number": 0,
                    }
                )
        elif ext == ".txt":
            content = self.load_txt(file_path)
            if content:
                documents.append(
                    {
                        "content": content,
                        "filename": filename,
                        "filepath": file_path,
                        "filetype": ext,
                        "page_number": 0,
                    }
                )
        else:
            print(f"不支持的文件格式: {ext}")

        return documents

    def load_all_documents(self) -> List[Dict[str, str]]:
        """加载数据目录下的所有文档"""
        if not os.path.exists(self.data_dir):
            print(f"数据目录不存在: {self.data_dir}")
            return None

        documents = []

        for root, dirs, files in os.walk(self.data_dir):
            for file in files:
                ext = os.path.splitext(file)[1].lower()
                if ext in self.supported_formats:
                    file_path = os.path.join(root, file)
                    print(f"正在加载: {file_path}")
                    doc_chunks = self.load_document(file_path)
                    if doc_chunks:
                        documents.extend(doc_chunks)

        return documents